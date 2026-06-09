#!/usr/bin/env python3
"""Build RangeGuard-Demo-Deck.pptx — 6-slide Google-Slides-ready deck (rebuild).

Structure (Session-15 rebuild):
  1 Title (logo-full + presenter + partner logos)
  2 The Solution (headline + 5 bullets + tagline)
  3 Economic Flywheel (two-row flow loop)
  4 Five Pillars
  5 Code Walkthrough (lifecycle + 2 amber callouts)
  6 Closing (standalone logo + bullets + beneficiary table + partner logos + links)

Logos are SVG but python-pptx embeds raster only, so docs/build_assets.py
pre-renders each to a navy-background PNG (seamless on the #0f1117 slide).
Run `python3 docs/build_assets.py` first, then this script.

Design system: bg #0f1117 · white #ffffff · accent #00d395 · slate #94a3b8 ·
amber #f59e0b · danger #ef4444 · card #1e2433 · Calibri · 16:9 widescreen.
"""
import os
from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ---------------------------------------------------------------- palette
BG     = RGBColor(0x0F, 0x11, 0x17)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
ACCENT = RGBColor(0x00, 0xD3, 0x95)
SLATE  = RGBColor(0x94, 0xA3, 0xB8)
AMBER  = RGBColor(0xF5, 0x9E, 0x0B)
DANGER = RGBColor(0xEF, 0x44, 0x44)
CARD   = RGBColor(0x1E, 0x24, 0x33)

HEAD_FONT = "Calibri"
BODY_FONT = "Calibri"
CODE_FONT = "Consolas"

EMU_IN = 914400
HERE = os.path.dirname(os.path.abspath(__file__))
PNG = os.path.join(HERE, "assets", "png")

prs = Presentation()
prs.slide_width  = Emu(int(13.333 * EMU_IN))
prs.slide_height = Emu(int(7.5 * EMU_IN))
BLANK = prs.slide_layouts[6]

_ASPECT = {}
def aspect(name):
    if name not in _ASPECT:
        w, h = Image.open(os.path.join(PNG, f"{name}.png")).size
        _ASPECT[name] = w / h
    return _ASPECT[name]


# ---------------------------------------------------------------- helpers
def add_slide():
    s = prs.slides.add_slide(BLANK)
    r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    r.fill.solid(); r.fill.fore_color.rgb = BG
    r.line.fill.background(); r.shadow.inherit = False
    sp = r._element; sp.getparent().remove(sp); s.shapes._spTree.insert(2, sp)
    return s


def _set_run(run, rd):
    run.text = rd["t"]
    f = run.font
    f.size = Pt(rd.get("size", 18)); f.bold = rd.get("bold", False)
    f.italic = rd.get("italic", False); f.name = rd.get("font", BODY_FONT)
    f.color.rgb = rd.get("color", WHITE)


def txt(slide, x, y, w, h, lines, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, wrap=True):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = wrap; tf.vertical_anchor = anchor
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    for i, para in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        if isinstance(para, dict):
            para = [para]
        meta = para[0] if para else {}
        if meta.get("space_before") is not None: p.space_before = Pt(meta["space_before"])
        if meta.get("space_after") is not None:  p.space_after = Pt(meta["space_after"])
        if meta.get("line_spacing") is not None: p.line_spacing = meta["line_spacing"]
        for rd in para:
            _set_run(p.add_run(), rd)
    return tb


def center(slide, x, y, w, h, paras, anchor=MSO_ANCHOR.MIDDLE):
    return txt(slide, x, y, w, h, paras, align=PP_ALIGN.CENTER, anchor=anchor)


def card(slide, x, y, w, h, fill=CARD, border=None, border_w=1.5, radius=0.08):
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                Inches(x), Inches(y), Inches(w), Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    if border is None: sh.line.fill.background()
    else: sh.line.color.rgb = border; sh.line.width = Pt(border_w)
    sh.shadow.inherit = False
    try: sh.adjustments[0] = radius
    except Exception: pass
    return sh


def line(slide, x, y, w, color=ACCENT, weight=2.5):
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                Inches(x), Inches(y), Inches(w), Inches(weight / 72.0))
    sh.fill.solid(); sh.fill.fore_color.rgb = color
    sh.line.fill.background(); sh.shadow.inherit = False
    return sh


def arrow_shape(slide, x, y, w, h, color=ACCENT, shape=MSO_SHAPE.RIGHT_ARROW):
    sh = slide.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = color
    sh.line.fill.background(); sh.shadow.inherit = False
    return sh


def img(slide, name, x, y, h):
    """Place PNG by height (inches), preserving aspect. Returns (x, y, w, h)."""
    w = h * aspect(name)
    pic = slide.shapes.add_picture(os.path.join(PNG, f"{name}.png"),
                                   Inches(x), Inches(y), height=Inches(h))
    pic.line.fill.background()   # ensure no faint picture border at the raster edge
    pic.shadow.inherit = False
    return x, y, w, h


def img_center(slide, name, cx, y, h):
    w = h * aspect(name)
    return img(slide, name, cx - w / 2, y, h)


def label(slide, text):
    txt(slide, 0.92, 0.46, 6.0, 0.4,
        [[{"t": text, "size": 13, "color": ACCENT, "bold": True, "font": HEAD_FONT}]])
    img(slide, "logo-icon", 0.5, 0.4, 0.3)   # subtle 24px-ish icon top-left


def notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text


def _text_w(s, size):
    return len(s) * 0.46 * size / 72.0


def hrow(slide, cy, items, gap=0.14, cx=6.6665):
    """Lay mixed text/image items in one horizontal row, centered at cx, middle at cy."""
    widths = []
    for it in items:
        if it["type"] == "text":
            widths.append(_text_w(it["t"], it.get("size", 13)) + 0.04)
        else:
            widths.append(it["h"] * aspect(it["name"]))
    total = sum(widths) + gap * (len(items) - 1)
    x = cx - total / 2
    for it, w in zip(items, widths):
        if it["type"] == "text":
            th = it.get("size", 13) / 72.0 * 1.8
            txt(slide, x, cy - th / 2, w + 0.1, th,
                [[{"t": it["t"], "size": it.get("size", 13),
                   "color": it.get("color", SLATE), "bold": it.get("bold", False),
                   "font": it.get("font", BODY_FONT)}]],
                align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE, wrap=False)
        else:
            img(slide, it["name"], x, cy - it["h"] / 2, it["h"])
        x += w + gap


# ================================================================ SLIDE 1 — Title
s = add_slide()
# Full lockup composed natively: shield icon + crisp Calibri name + tagline
img_center(s, "logo-icon", 6.6665, 1.0, 1.5)
center(s, 1.0, 2.72, 11.33, 0.95,
       [[{"t": "RangeGuard", "size": 48, "color": WHITE, "bold": True, "font": HEAD_FONT}]],
       anchor=MSO_ANCHOR.TOP)
center(s, 1.0, 3.78, 11.33, 0.5,
       [[{"t": "Protect your liquidity. Guard your range.", "size": 20, "color": SLATE}]],
       anchor=MSO_ANCHOR.TOP)
center(s, 1.0, 4.62, 11.33, 0.95, [
    [{"t": "Gary Kocsis", "size": 22, "color": WHITE, "bold": True, "space_after": 3}],
    [{"t": "Uniswap Hook Incubator", "size": 16, "color": SLATE}],
], anchor=MSO_ANCHOR.TOP)
hrow(s, 6.35, [
    {"type": "text", "t": "Built on", "size": 14, "color": SLATE},
    {"type": "img", "name": "uniswap-logo", "h": 0.34},
    {"type": "text", "t": "·", "size": 18, "color": SLATE},
    {"type": "text", "t": "Powered by", "size": 14, "color": SLATE},
    {"type": "img", "name": "reactive-logo", "h": 0.30},
], gap=0.22)
notes(s,
"Hello everyone, my name is Gary Kocsis, and I'm excited to present RangeGuard — Protect your "
"liquidity. Guard your range. As many of you know, impermanent loss is the single biggest reason "
"liquidity providers leave AMMs. Even with fee rewards, most LPs lose money when prices move — "
"and there's no native, on-chain way to protect them. Until now.")

# ================================================================ SLIDE 2 — The Solution
s = add_slide()
label(s, "THE SOLUTION")
center(s, 1.0, 1.05, 11.33, 1.2, [
    [{"t": "RangeGuard turns impermanent loss", "size": 30, "color": WHITE, "bold": True, "font": HEAD_FONT}],
    [{"t": "into an earned, capped, on-chain claim.", "size": 30, "color": WHITE, "bold": True, "font": HEAD_FONT}],
], anchor=MSO_ANCHOR.TOP)
sol_bullets = [
    ("Coverage accrues only while in range —", "exposure-weighted, not time-weighted"),
    ("Day-count convention (Actual/365 Fixed) —", "predictable, auditable, comparable across pools"),
    ("Buffer funded by dynamic fee skimming —", "sustainable, no external subsidies"),
    ("Automatic payout on withdrawal —", "capped by earned coverage, buffer health, and protocol params"),
    ("Reactive Network integration —", "range transitions detected cross-chain, no keepers"),
]
y = 2.55
for top, sub in sol_bullets:
    txt(s, 1.5, y, 10.5, 0.78, [
        [{"t": "▸  ", "size": 17, "color": ACCENT, "bold": True},
         {"t": top, "size": 17, "color": WHITE, "bold": True}],
        [{"t": "     " + sub, "size": 15, "color": SLATE, "space_before": 1}],
    ])
    y += 0.82
center(s, 1.0, 6.85, 11.33, 0.45,
       [[{"t": "“Protect your liquidity. Guard your range.”",
          "size": 19, "color": ACCENT, "bold": True, "italic": True}]])
notes(s,
"RangeGuard brings native, transparent IL coverage directly into a Uniswap v4 pool. When an LP "
"provides liquidity, they start earning coverage — but only while their position is in range and "
"exposed to IL risk. Coverage accrues using a day-count convention — just like traditional "
"finance — so LPs can see exactly how much protection they're earning and compare it across "
"pools. The coverage buffer is funded by a small configurable portion of trading fees via "
"Uniswap v4's dynamic fee system — sustainable, no external subsidies. When an LP withdraws, the "
"hook automatically calculates their impermanent loss and pays out a capped reimbursement — "
"bounded by what they've earned, the buffer's health, and protocol parameters. And the system "
"runs autonomously. RangeGuard integrates with the Reactive Network — a cross-chain automation "
"layer running on Lasna — so range transitions are detected and recorded automatically. No "
"keeper bots. No off-chain infrastructure.")

# ================================================================ SLIDE 3 — Economic Flywheel
s = add_slide()
label(s, "THE SOLUTION")
center(s, 1.0, 0.95, 11.33, 0.6,
       [[{"t": "Self-Funding by Design", "size": 30, "color": WHITE, "bold": True, "font": HEAD_FONT}]],
       anchor=MSO_ANCHOR.TOP)
# two-row flow: 3 columns
bw, bh = 3.35, 1.05
c = [1.05, 4.99, 8.93]   # column x's
r1y, r2y = 2.25, 4.35
row1 = ["LPs provide\nliquidity", "Traders use\nthe pool", "Swaps generate\nfees"]
row2 = ["LPs are\nprotected", "Buffer pays\ncapped IL", "Buffer slice\nskimmed"]
row2_hl = [False, True, True]
def fbox(x, y, text, hl):
    card(s, x, y, bw, bh, fill=CARD, border=(ACCENT if hl else None), border_w=2)
    center(s, x, y, bw, bh,
           [[{"t": ln, "size": 15, "color": WHITE, "bold": hl}] for ln in text.split("\n")])
for i in range(3):
    fbox(c[i], r1y, row1[i], False)
    fbox(c[i], r2y, row2[i], row2_hl[i])
# row1 right arrows (between cols)
for i in range(2):
    ax = c[i] + bw; aw = c[i+1] - (c[i] + bw)
    arrow_shape(s, ax + 0.05, r1y + bh/2 - 0.18, aw - 0.1, 0.36, color=ACCENT, shape=MSO_SHAPE.RIGHT_ARROW)
# row2 left arrows (between cols, pointing left)
for i in range(2):
    ax = c[i] + bw; aw = c[i+1] - (c[i] + bw)
    arrow_shape(s, ax + 0.05, r2y + bh/2 - 0.18, aw - 0.1, 0.36, color=ACCENT, shape=MSO_SHAPE.LEFT_ARROW)
# down arrow on right column (row1 -> row2)
arrow_shape(s, c[2] + bw/2 - 0.18, r1y + bh + 0.06, 0.36, r2y - (r1y + bh) - 0.12,
            color=ACCENT, shape=MSO_SHAPE.DOWN_ARROW)
# up arrow on left column (row2 -> row1)
arrow_shape(s, c[0] + bw/2 - 0.18, r1y + bh + 0.06, 0.36, r2y - (r1y + bh) - 0.12,
            color=ACCENT, shape=MSO_SHAPE.UP_ARROW)
center(s, 1.0, 6.7, 11.33, 0.5,
       [[{"t": "Swap activity funds the buffer that protects LPs.",
          "size": 19, "color": ACCENT, "bold": True}]])
notes(s,
"The economics are self-sustaining. LPs provide liquidity, traders use the pool, swaps generate "
"fees, a buffer slice is skimmed that funds the coverage buffer, and the buffer pays capped IL "
"claims. The buffer grows from the same swap activity that LP liquidity enables. No external "
"capital required.")

# ================================================================ SLIDE 4 — Five Pillars
s = add_slide()
label(s, "THE SOLUTION")
center(s, 1.0, 0.95, 11.33, 0.6,
       [[{"t": "How RangeGuard Works: Five Pillars", "size": 29, "color": WHITE, "bold": True, "font": HEAD_FONT}]],
       anchor=MSO_ANCHOR.TOP)
pillars = [
    ("1", "Accrual Gating", ["Coverage accrues only while the position is in range"], False),
    ("2", "Buffer Funding", ["A portion of every swap fee funds the coverage buffer"], False),
    ("3", "Claim Settlement", ["IL computed and paid automatically on withdrawal",
                               "Payout = min(covered IL, earned coverage, buffer cap)"], False),
    ("4", "LP Transparency  ★", ["A verifiable, day-by-day coverage report —",
                                 "built entirely from on-chain events"], True),
    ("5", "Pool Configuration", ["Immutable parameters set once at pool initialization"], False),
]
y = 1.8
for num, title, body, hl in pillars:
    rowh = 0.74 + 0.26 * (len(body) - 1)
    if hl:
        card(s, 0.85, y - 0.08, 11.63, rowh + 0.12, fill=CARD, border=ACCENT, border_w=2)
    txt(s, 1.1, y, 0.7, rowh, [[{"t": num, "size": 25, "color": ACCENT, "bold": True}]])
    paras = [[{"t": title, "size": 19, "color": WHITE, "bold": True, "space_after": 2}]]
    for b in body:
        paras.append([{"t": b, "size": 14.5, "color": SLATE, "space_after": 1}])
    txt(s, 1.8, y, 10.5, rowh, paras)
    y += rowh + 0.16
center(s, 1.0, 6.85, 11.33, 0.45,
       [[{"t": "Every pillar is enforced on-chain. No off-chain assumptions.",
          "size": 17, "color": ACCENT, "bold": True}]])
notes(s,
"RangeGuard is built on five pillars. Accrual gating — coverage only earns while in range. Buffer "
"funding — every swap contributes automatically. Claim settlement — IL is computed and paid in "
"the same withdrawal transaction, capped by three limits: covered IL, earned coverage, and "
"buffer capacity. LP transparency — the key differentiator — a verifiable day-by-day coverage "
"statement from pure on-chain events. And pool configuration — immutable parameters, so LPs "
"always know what they signed up for.\n\n"
"Verbal transition: Let me show you the core functions that make this work.")

# ================================================================ SLIDE 5 — Code Walkthrough
s = add_slide()
label(s, "UNDER THE HOOD")
center(s, 1.0, 0.95, 11.33, 0.6,
       [[{"t": "How Coverage Becomes a Claim", "size": 29, "color": WHITE, "bold": True, "font": HEAD_FONT}]],
       anchor=MSO_ANCHOR.TOP)
steps = ["LP Deposits", "_accrue()", "afterSwap()", "_computeIL()", "_computePayout()", "ClaimSettled"]
n = len(steps); fx = 0.55; fw = 1.78; fgap = (12.78 - 0.55 - n * fw) / (n - 1)
fy = 2.2; fh = 0.95
for i, lab in enumerate(steps):
    x = fx + i * (fw + fgap)
    card(s, x, fy, fw, fh, fill=CARD, border=ACCENT, border_w=2)   # all same color
    center(s, x, fy, fw, fh,
           [[{"t": lab, "size": 11.5, "color": WHITE, "bold": True, "font": CODE_FONT}]])
    if i < n - 1:
        ax = x + fw + (fgap - 0.4) / 2
        txt(s, ax, fy + fh/2 - 0.25, 0.5, 0.5,
            [[{"t": "→", "size": 24, "color": ACCENT, "bold": True}]],
            align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
card(s, 0.85, 3.95, 5.6, 1.95, fill=CARD, border=AMBER, border_w=2)
txt(s, 1.15, 4.2, 5.05, 1.5, [
    [{"t": "⚡  Swaps never iterate positions", "size": 16, "color": AMBER, "bold": True, "space_after": 6}],
    [{"t": "Accrual is lazy, position-specific", "size": 15, "color": WHITE, "space_after": 4}],
    [{"t": "O(1) gas on every swap", "size": 15, "color": WHITE, "bold": True}],
])
card(s, 6.75, 3.95, 5.7, 1.95, fill=CARD, border=AMBER, border_w=2)
txt(s, 7.05, 4.2, 5.2, 1.55, [
    [{"t": "⚡  Actual/365 Fixed day-count", "size": 16, "color": AMBER, "bold": True, "space_after": 6}],
    [{"t": "Coverage = Notional × APR × (days ÷ 365)", "size": 14, "color": WHITE, "font": CODE_FONT, "space_after": 4}],
    [{"t": "Same convention as fixed-income finance", "size": 15, "color": WHITE}],
])
center(s, 1.0, 6.4, 11.33, 0.5,
       [[{"t": "One lifecycle. Fully on-chain. No off-chain assumptions.",
          "size": 18, "color": ACCENT, "bold": True}]])
notes(s,
"Before the demo, here's the core code path. _accrue advances coverage lazily — only while in "
"range, using Actual/365 Fixed day-count math. _computePayout applies three caps in order. The "
"day-count convention turns coverage from a black-box reward into an auditable financial "
"accrual. Let me show you.\n\n"
"IDE narration — _accrue: range gate first — if the tick is outside the LP's bounds, delta is "
"zero. If in range — year fraction from day-count basis, multiplied by entry notional and APR. "
"Fifteen days in range earns exactly 15/365 of the annual coverage amount. _computePayout: three "
"caps in order — IL cap, earned coverage cap, buffer cap. Minimum of all three. Limiting factor "
"is recorded so the LP knows exactly why they received what they received.")

# ================================================================ SLIDE 6 — Closing
s = add_slide()
# Standalone lockup composed natively: shield icon + crisp Calibri name, centered
_name_w = 2.05  # generous, avoids wrap of bold 27pt "RangeGuard"
_icon_h = 0.52
_grp = _icon_h + 0.12 + _name_w
_gx = 6.6665 - _grp / 2
img(s, "logo-icon", _gx, 0.29, _icon_h)
txt(s, _gx + _icon_h + 0.12, 0.3, _name_w, 0.5,
    [[{"t": "RangeGuard", "size": 27, "color": WHITE, "bold": True, "font": HEAD_FONT}]],
    align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE, wrap=False)
center(s, 1.0, 0.95, 11.33, 0.55,
       [[{"t": "“Protect your liquidity. Guard your range.”",
          "size": 22, "color": ACCENT, "bold": True, "italic": True}]])
fbul = [
    "Native IL coverage — built into the pool, not bolted on",
    "Self-funding buffer — swap fees cover the claims",
    "Capped payouts — actuarially sound, buffer protected",
    "Fully auditable — every accrual verifiable on-chain",
]
txt(s, 2.3, 1.65, 9.0, 1.25,
    [[{"t": "✓  ", "size": 16, "color": ACCENT, "bold": True},
      {"t": b, "size": 16, "color": WHITE}] for b in
     [f for f in fbul]],
    )
# beneficiary table
ty = 3.0
card(s, 1.7, ty, 9.93, 1.6, fill=CARD)
bens = [
    ("Who Benefits", "Why It Matters", True),
    ("Passive LPs", "Downside protection without complex hedging", False),
    ("Protocols / DAOs", "Deeper liquidity without reliance on emissions", False),
    ("Uniswap Ecosystem", "Durable liquidity, tighter spreads, better execution", False),
]
by0 = ty + 0.12; brh = 0.36
for j, (a, b, hdr) in enumerate(bens):
    yy = by0 + j * brh
    txt(s, 2.05, yy, 3.3, brh,
        [[{"t": a, "size": 13.5, "color": (ACCENT if hdr else WHITE), "bold": True}]],
        anchor=MSO_ANCHOR.MIDDLE)
    txt(s, 5.4, yy, 6.0, brh,
        [[{"t": b, "size": 13, "color": SLATE}]], anchor=MSO_ANCHOR.MIDDLE)
    if hdr:
        line(s, 2.05, yy + brh - 0.02, 9.2, color=SLATE, weight=1)
center(s, 1.0, 4.78, 11.33, 0.45,
       [[{"t": "Earned over time. Funded by swaps. Settled on-chain.",
          "size": 19, "color": ACCENT, "bold": True}]])
line(s, 5.165, 5.3, 3.0, color=SLATE, weight=1)
center(s, 1.0, 5.42, 11.33, 0.4,
       [[{"t": "🔗  range-guard.vercel.app      ·      github.com/garykocsis/RangeGuard",
          "size": 14, "color": WHITE, "bold": True}]])
hrow(s, 6.15, [
    {"type": "text", "t": "Built on", "size": 13, "color": SLATE},
    {"type": "img", "name": "uniswap-logo", "h": 0.32},
    {"type": "text", "t": "·", "size": 16, "color": SLATE},
    {"type": "text", "t": "Powered by", "size": 13, "color": SLATE},
    {"type": "img", "name": "reactive-logo", "h": 0.28},
], gap=0.2)
center(s, 1.0, 6.95, 11.33, 0.35,
       [[{"t": "292 tests passing   ·   Fuzz tested   ·   Invariant tested",
          "size": 12, "color": SLATE}]])
notes(s,
"RangeGuard makes providing liquidity safer, more transparent, and more attractive — helping "
"Uniswap pools retain and grow their LP base. The primary beneficiaries are passive LPs who get "
"downside protection without complex hedging. But the impact extends to protocols that need "
"sticky liquidity without emissions, and to the broader ecosystem through deeper pools and "
"better execution. Earned over time. Funded by swaps. Settled on-chain. Live dashboard at "
"range-guard.vercel.app. Full source and 292 passing tests on GitHub. Thank you.")

# ---------------------------------------------------------------- save
out = os.path.join(HERE, "RangeGuard-Demo-Deck.pptx")
prs.save(out)
print("saved", out, "slides:", len(prs.slides._sldIdLst))
